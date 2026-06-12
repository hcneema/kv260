from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
import os

prs = Presentation()
prs.slide_width = Inches(13.33)
prs.slide_height = Inches(7.5)

DARK_BG    = RGBColor(0x1E, 0x1E, 0x2E)
AMD_RED    = RGBColor(0xED, 0x1C, 0x24)
WHITE      = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY = RGBColor(0xCC, 0xCC, 0xCC)
GREEN      = RGBColor(0x00, 0xC8, 0x53)
YELLOW     = RGBColor(0xFF, 0xD6, 0x00)
BLUE       = RGBColor(0x29, 0xB6, 0xF6)
ORANGE     = RGBColor(0xFF, 0x6D, 0x00)

blank_layout = prs.slide_layouts[6]

def add_slide():
    slide = prs.slides.add_slide(blank_layout)
    bg = slide.background.fill
    bg.solid()
    bg.fore_color.rgb = DARK_BG
    return slide

def add_text(slide, text, x, y, w, h, size=18, bold=False, color=WHITE, align=PP_ALIGN.LEFT, wrap=True):
    txBox = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = txBox.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    return txBox

def add_rect(slide, x, y, w, h, fill_color, line_color=None):
    shape = slide.shapes.add_shape(1, Inches(x), Inches(y), Inches(w), Inches(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if line_color:
        shape.line.color.rgb = line_color
    else:
        shape.line.fill.background()
    return shape

def add_title_bar(slide, title, subtitle=None):
    add_rect(slide, 0, 0, 13.33, 1.3, AMD_RED)
    add_text(slide, title, 0.3, 0.1, 12, 0.7, size=32, bold=True, color=WHITE)
    if subtitle:
        add_text(slide, subtitle, 0.3, 0.8, 12, 0.5, size=16, color=RGBColor(0xFF,0xCC,0xCC))

# ─── SLIDE 1: Title ───────────────────────────────────────────────
slide = add_slide()
add_rect(slide, 0, 0, 13.33, 7.5, DARK_BG)
add_rect(slide, 0, 0, 13.33, 0.08, AMD_RED)
add_rect(slide, 0, 7.42, 13.33, 0.08, AMD_RED)

add_text(slide, 'AMD Kria KV260', 1, 1.5, 11, 1.2, size=52, bold=True, color=AMD_RED, align=PP_ALIGN.CENTER)
add_text(slide, 'Vision Robot Project', 1, 2.7, 11, 1, size=40, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
add_text(slide, 'From Zero to Face-Detection + Servo Control in One Day', 1, 3.8, 11, 0.6, size=20, color=LIGHT_GRAY, align=PP_ALIGN.CENTER)
add_text(slide, 'KV260  ·  Elegoo Mega 2560  ·  Ubuntu 24.04  ·  Python  ·  OpenCV', 1, 5.5, 11, 0.5, size=16, color=RGBColor(0x88,0x88,0x88), align=PP_ALIGN.CENTER)
add_text(slide, 'June 2026', 1, 6.5, 11, 0.5, size=14, color=RGBColor(0x66,0x66,0x66), align=PP_ALIGN.CENTER)

# ─── SLIDE 2: Hardware Overview ───────────────────────────────────
slide = add_slide()
add_title_bar(slide, '🔧 Hardware Stack', 'All components and how they connect')

items = [
    (AMD_RED,  '🧠  KV260',         'AMD Kria KV260 — ARM Cortex-A53 + FPGA DPU  |  IP: 192.168.68.60  |  Ubuntu 24.04'),
    (BLUE,     '⚡  Elegoo Mega',   'Arduino Mega 2560 compatible — runs motor/servo control sketch'),
    (ORANGE,   '🔌  L298P Shield',  '4-channel motor driver — plugs ON TOP of Mega  |  2A per channel  |  screw terminals'),
    (GREEN,    '⚙️  TT DC Motors',  '2x TT Gear Motors (1:48 ratio, 3-6V)  +  wheels — press fit onto shaft'),
    (YELLOW,   '📷  USB Webcam',    'Logitech UVC  |  /dev/video0  |  640×480  |  live inference at 12fps'),
    (WHITE,    '🎛️  Elegoo Kit',    'SG90 Servo  ·  HC-SR04 Ultrasonic  ·  IR Remote  ·  GY-521 IMU — all reused!'),
]

for i, (color, title, desc) in enumerate(items):
    y = 1.5 + i * 0.9
    add_rect(slide, 0.3, y, 0.08, 0.55, color)
    add_text(slide, title, 0.55, y, 3, 0.55, size=15, bold=True, color=color)
    add_text(slide, desc,  3.7,  y, 9, 0.55, size=13, color=LIGHT_GRAY)

# ─── SLIDE 3: System Architecture ────────────────────────────────
slide = add_slide()
add_title_bar(slide, '🏗️ System Architecture', 'Full signal flow from camera to wheels')

boxes = [
    (0.4,  3.2, 2.2, 1.1, BLUE,   '📷 USB Webcam\n/dev/video0\n640×480'),
    (3.2,  3.2, 2.5, 1.1, AMD_RED,'🧠 KV260\nPython + OpenCV\nFace Detection'),
    (6.2,  3.2, 2.5, 1.1, ORANGE, '⚡ Elegoo Mega\nArduino Sketch\n9600 baud serial'),
    (9.2,  2.6, 1.8, 1.0, GREEN,  '⚙️ Left\nMotor'),
    (9.2,  3.9, 1.8, 1.0, GREEN,  '⚙️ Right\nMotor'),
    (11.3, 2.6, 1.7, 1.0, YELLOW, '🛞 Left\nWheel'),
    (11.3, 3.9, 1.7, 1.0, YELLOW, '🛞 Right\nWheel'),
]
arrows = [
    (2.6, 3.75, 3.2, 3.75),
    (5.7, 3.75, 6.2, 3.75),
    (8.7, 3.1,  9.2, 3.1),
    (8.7, 4.4,  9.2, 4.4),
    (11.0,3.1,  11.3,3.1),
    (11.0,4.4,  11.3,4.4),
]
labels = [
    (2.6, 3.4,  'frames'),
    (5.7, 3.4,  'serial\n/dev/ttyACM0'),
    (8.7, 2.8,  'PWM'),
    (8.7, 4.1,  'PWM'),
]

for (x,y,w,h,col,txt) in boxes:
    add_rect(slide, x, y, w, h, col)
    add_text(slide, txt, x+0.05, y+0.05, w-0.1, h-0.1, size=12, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

for (x1,y1,x2,y2) in arrows:
    from pptx.util import Inches as I
    line = slide.shapes.add_connector(1, I(x1), I(y1), I(x2), I(y2))
    line.line.color.rgb = WHITE
    line.line.width = Pt(2)

for (x,y,txt) in labels:
    add_text(slide, txt, x, y, 1.5, 0.4, size=10, color=LIGHT_GRAY)

add_rect(slide, 8.5, 2.0, 2.5, 0.6, RGBColor(0x33,0x33,0x44))
add_text(slide, 'L298P Motor Shield', 8.5, 2.05, 2.5, 0.5, size=11, bold=True, color=ORANGE, align=PP_ALIGN.CENTER)

add_text(slide, '🔌 USB Cable', 2.6, 4.15, 1.5, 0.35, size=10, color=LIGHT_GRAY)
add_text(slide, '🔋 AA Batteries → Shield → Motors\n🔌 Wall Adapter → KV260', 0.3, 6.0, 8, 0.8, size=12, color=LIGHT_GRAY)

# ─── SLIDE 4: What's Working Today ───────────────────────────────
slide = add_slide()
add_title_bar(slide, '✅ Proven Working — Day 1', 'Everything achieved in a single session')

achievements = [
    (GREEN,  '✅', 'SSH from Windows PowerShell → KV260',          'Fixed Sysnative path issue, added alias to PS profile'),
    (GREEN,  '✅', 'Serial link: KV260 ↔ Elegoo Mega',            '/dev/ttyACM0 at 9600 baud — Hello from Elegoo confirmed'),
    (GREEN,  '✅', 'Servo control from KV260',                     'Python pyserial sends angles 0→90→180, servo moves on command'),
    (GREEN,  '✅', 'USB Webcam detected on KV260',                 'Logitech UVC on /dev/video0, 640×480 resolution'),
    (GREEN,  '✅', 'Live object detection — MobileNetV2',          '~83ms/frame (12fps) on ARM CPU, detects couch/desk/person'),
    (GREEN,  '✅', 'Face detection — Haar Cascade',                'OpenCV built-in, reliable at 1-2m distance'),
    (GREEN,  '✅', 'FULL PIPELINE: Face → Serial → Servo',         'Webcam sees face → KV260 → Mega → servo sweeps 0→90→180→90→0'),
]

for i, (color, icon, title, desc) in enumerate(achievements):
    y = 1.5 + i * 0.77
    add_rect(slide, 0.3, y+0.1, 0.5, 0.5, color)
    add_text(slide, icon,  0.3,  y+0.05, 0.6, 0.55, size=18, align=PP_ALIGN.CENTER)
    add_text(slide, title, 1.0,  y+0.05, 5.5, 0.4,  size=14, bold=True, color=WHITE)
    add_text(slide, desc,  1.0,  y+0.42, 11,  0.35, size=11, color=LIGHT_GRAY)

# ─── SLIDE 5: Development Roadmap ────────────────────────────────
slide = add_slide()
add_title_bar(slide, '🗺️ Development Roadmap', 'CPU-first approach — DPU when ready')

phases = [
    (0,  GREEN,  '✅ Phase 0',  'DONE',    'SSH · Serial · Servo · Webcam · Face Detection · Servo Sweep'),
    (1,  BLUE,   '🔜 Phase 1', 'PENDING', 'Motor control — KV260 drives 2x DC motors via L298P shield'),
    (2,  BLUE,   '🔜 Phase 2', 'PENDING', 'IR Remote control — drive robot with Elegoo IR remote'),
    (3,  BLUE,   '🔜 Phase 3', 'PENDING', 'Obstacle avoidance — HC-SR04 ultrasonic sensor'),
    (4,  BLUE,   '🔜 Phase 4', 'PENDING', 'Live webcam object detection on moving robot (CPU, 12fps)'),
    (5,  ORANGE, '🚀 Phase 5', 'FUTURE',  'YOLO on DPU — real-time 30fps AI vision (needs Ubuntu 22.04)'),
]

for i, (_, color, phase, status, desc) in enumerate(phases):
    y = 1.5 + i * 0.87
    add_rect(slide, 0.3,  y, 2.2, 0.65, color)
    add_text(slide, phase,  0.35, y+0.08, 2.1, 0.5, size=14, bold=True, color=DARK_BG, align=PP_ALIGN.CENTER)
    add_rect(slide, 2.7,  y, 1.5, 0.65, RGBColor(0x33,0x33,0x44))
    add_text(slide, status, 2.75, y+0.08, 1.4, 0.5, size=12, bold=True, color=color, align=PP_ALIGN.CENTER)
    add_text(slide, desc,   4.4,  y+0.1,  8.6, 0.5, size=13, color=LIGHT_GRAY)

add_text(slide, '* Phases 1-4 run on Ubuntu 24.04 CPU only — no new microSD needed until Phase 5',
         0.3, 6.9, 12.5, 0.4, size=11, color=RGBColor(0x88,0x88,0x88))

# ─── SLIDE 6: Shopping List ───────────────────────────────────────
slide = add_slide()
add_title_bar(slide, '🛒 Shopping List', 'Everything needed — under $50 to get motors spinning')

add_text(slide, 'Microcenter Santa Clara — 5201 Stevens Creek Blvd', 0.5, 1.5, 12, 0.5, size=16, bold=True, color=YELLOW)
mc_items = [
    ('L298P 4-Channel Motor Drive Shield', '~$30', 'Confirmed in stock'),
    ('MicroSD Card 32GB+',                 '~$10', 'For Ubuntu 22.04 flash (Phase 5)'),
]
for i, (item, price, note) in enumerate(mc_items):
    y = 2.1 + i * 0.65
    add_rect(slide, 0.5, y, 7.5, 0.5, RGBColor(0x2A,0x2A,0x3A))
    add_text(slide, item,  0.65, y+0.07, 5.5, 0.38, size=13, color=WHITE)
    add_text(slide, price, 6.2,  y+0.07, 1.0, 0.38, size=13, bold=True, color=GREEN, align=PP_ALIGN.CENTER)
    add_text(slide, note,  7.3,  y+0.07, 3.0, 0.38, size=11, color=LIGHT_GRAY)

add_text(slide, 'Amazon', 0.5, 3.6, 12, 0.5, size=16, bold=True, color=ORANGE)
amz_items = [
    ('DIYmall 2x TT DC Gear Motor + 2 Wheels', '~$10', 'Initial desk test — no chassis needed yet'),
    ('XiaoR Geek Aluminum Tank Chassis 2WD',   '~$40', 'Future — after motors verified working'),
    ('12V USB-C PD Power Bank',                 '~$40', 'Future — for untethered robot operation'),
    ('Raspberry Pi Camera v2',                  '~$25', 'Future — for DPU-accelerated vision'),
]
for i, (item, price, note) in enumerate(amz_items):
    y = 4.2 + i * 0.65
    add_rect(slide, 0.5, y, 7.5, 0.5, RGBColor(0x2A,0x2A,0x3A))
    add_text(slide, item,  0.65, y+0.07, 5.5, 0.38, size=13, color=WHITE)
    add_text(slide, price, 6.2,  y+0.07, 1.0, 0.38, size=13, bold=True, color=GREEN, align=PP_ALIGN.CENTER)
    add_text(slide, note,  7.3,  y+0.07, 5.5, 0.38, size=11, color=LIGHT_GRAY)

add_text(slide, 'Immediate total: ~$50   |   Full robot total: ~$155', 0.5, 6.9, 12, 0.4, size=13, bold=True, color=YELLOW)

# ─── SLIDE 7: Face Detection Demo ────────────────────────────────
slide = add_slide()
add_title_bar(slide, '🎯 Face Detection + Servo Demo', 'Working today — no DPU needed!')

add_text(slide, 'How it works:', 0.5, 1.5, 12, 0.5, size=18, bold=True, color=WHITE)

steps = [
    (BLUE,   '1', 'Webcam captures frame',            '/dev/video0  →  640×480 image every 0.5s'),
    (BLUE,   '2', 'Convert to grayscale',             'cv2.cvtColor(frame, COLOR_BGR2GRAY)'),
    (AMD_RED,'3', 'Haar Cascade face detection',      'haarcascade_frontalface_default.xml  |  minNeighbors=3  |  minSize=30×30'),
    (GREEN,  '4', 'Face found?',                      'YES → trigger servo sweep  |  NO → keep scanning silently'),
    (ORANGE, '5', 'Send angle commands to Mega',       'pyserial → /dev/ttyACM0  |  0 → 90 → 180 → 90 → 0'),
    (YELLOW, '6', 'Mega moves servo',                  'servo_test.ino  |  Serial.parseInt()  |  myServo.write(angle)'),
]

for i, (color, num, title, desc) in enumerate(steps):
    y = 2.1 + i * 0.77
    add_rect(slide, 0.4, y+0.05, 0.45, 0.5, color)
    add_text(slide, num,   0.4,  y+0.08, 0.45, 0.45, size=18, bold=True, color=DARK_BG, align=PP_ALIGN.CENTER)
    add_text(slide, title, 1.05, y+0.05, 4.5,  0.4,  size=14, bold=True, color=color)
    add_text(slide, desc,  1.05, y+0.42, 11.5, 0.35, size=11, color=LIGHT_GRAY)

add_rect(slide, 0.4, 6.75, 12.4, 0.45, RGBColor(0x00,0x40,0x00))
add_text(slide, '✅  Result: Face detected → servo sweeps 0→90→180→90→0 every 10 seconds. No face → stays still.',
         0.5, 6.78, 12, 0.38, size=12, bold=True, color=GREEN)

# ─── SLIDE 8: Ultimate Goal ───────────────────────────────────────
slide = add_slide()
add_title_bar(slide, '🚀 Ultimate Goal', 'Autonomous vision robot powered by AMD DPU')

add_text(slide, 'The Final Robot', 0.5, 1.5, 12, 0.6, size=24, bold=True, color=WHITE)

goal_items = [
    (AMD_RED, '👁️  Sees',     'Camera captures live video → KV260 DPU runs YOLO at 30fps'),
    (BLUE,    '🧠  Thinks',   'Detects: person / obstacle / object — decides what to do'),
    (GREEN,   '⚙️  Acts',     'Sends motor commands → Mega → wheels move accordingly'),
    (ORANGE,  '🔄  Reacts',   'Follow person · avoid obstacle · drive toward target object'),
    (YELLOW,  '📡  Control',  'IR remote override · ROS2 navigation · autonomous exploration'),
]

for i, (color, icon_title, desc) in enumerate(goal_items):
    y = 2.3 + i * 0.85
    add_rect(slide, 0.4, y, 3.5, 0.65, color)
    add_text(slide, icon_title, 0.5, y+0.1, 3.3, 0.5, size=16, bold=True, color=DARK_BG)
    add_text(slide, desc, 4.1, y+0.12, 8.8, 0.45, size=13, color=LIGHT_GRAY)

add_rect(slide, 0.4, 6.6, 12.4, 0.65, RGBColor(0x1A,0x1A,0x2A))
add_text(slide, '💡  Think of it as a mini self-driving car — camera eyes, FPGA brain, Arduino muscles.',
         0.6, 6.65, 12, 0.55, size=14, bold=True, color=WHITE)

pptx_path = 'C:/Users/hemn/kv260/robotics/KV260_Robot_Project.pptx'
prs.save(pptx_path)
print('Saved:', pptx_path)
