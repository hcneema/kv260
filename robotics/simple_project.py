from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

prs = Presentation()
prs.slide_width = Inches(13.33)
prs.slide_height = Inches(7.5)

DARK_BG    = RGBColor(0x1E, 0x1E, 0x2E)
AMD_RED    = RGBColor(0xED, 0x1C, 0x24)
WHITE      = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY = RGBColor(0xCC, 0xCC, 0xCC)
GREEN      = RGBColor(0x00, 0xC8, 0x53)
BLUE       = RGBColor(0x29, 0xB6, 0xF6)
ORANGE     = RGBColor(0xFF, 0x6D, 0x00)
YELLOW     = RGBColor(0xFF, 0xD6, 0x00)

blank_layout = prs.slide_layouts[6]

def add_slide():
    slide = prs.slides.add_slide(blank_layout)
    bg = slide.background.fill
    bg.solid()
    bg.fore_color.rgb = DARK_BG
    return slide

def add_text(slide, text, x, y, w, h, size=18, bold=False, color=WHITE, align=PP_ALIGN.LEFT):
    txBox = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    return txBox

def add_rect(slide, x, y, w, h, fill_color, line_color=None):
    shape = slide.shapes.add_shape(1, Inches(x), Inches(y), Inches(w), Inches(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if line_color:
        shape.line.color.rgb = line_color
    else:
        shape.line.fill.background()
    return shape

def add_title_bar(slide, title):
    add_rect(slide, 0, 0, 13.33, 1.2, AMD_RED)
    add_text(slide, title, 0.4, 0.2, 12, 0.8, size=34, bold=True, color=WHITE)

# ─── SLIDE 1: Title ───────────────────────────────────────────────
slide = add_slide()
add_rect(slide, 0, 0, 13.33, 0.1, AMD_RED)
add_rect(slide, 0, 7.4, 13.33, 0.1, AMD_RED)
add_text(slide, 'Face Detection → Servo Control', 1, 1.8, 11, 1.2, size=44, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
add_text(slide, 'Using AMD Kria KV260 + Elegoo Mega 2560', 1, 3.1, 11, 0.7, size=24, color=LIGHT_GRAY, align=PP_ALIGN.CENTER)
add_text(slide, 'A working demo — built in one day', 1, 4.5, 11, 0.6, size=18, color=RGBColor(0x88,0x88,0x88), align=PP_ALIGN.CENTER)

# ─── SLIDE 2: The Pipeline (simple) ───────────────────────────────
slide = add_slide()
add_title_bar(slide, '📡 The Pipeline')

# boxes
components = [
    (0.4,  3.0, 2.0, 1.2, BLUE,   '📷\nUSB Webcam'),
    (3.2,  3.0, 2.5, 1.2, AMD_RED,'🧠\nKV260'),
    (6.4,  3.0, 2.5, 1.2, AMD_RED,'👤\nFace Detection'),
    (9.6,  3.0, 2.5, 1.2, ORANGE, '⚡\nElegoo Mega'),
    (12.3, 3.0, 0.9, 1.2, GREEN,  '🔄\nServo'),
]

for (x, y, w, h, col, txt) in components:
    add_rect(slide, x, y, w, h, col)
    add_text(slide, txt, x+0.05, y+0.1, w-0.1, h-0.15, size=14, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

# arrows
for x in [2.4, 5.7, 8.9, 12.1]:
    add_text(slide, '→', x, 3.35, 0.5, 0.5, size=24, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

# labels under boxes
labels = [
    (0.4,  4.35, 2.0, '/dev/video0\n640×480'),
    (3.2,  4.35, 2.5, 'Ubuntu 24.04\nPython + OpenCV'),
    (6.4,  4.35, 2.5, 'Haar Cascade\n~83ms per frame'),
    (9.6,  4.35, 2.5, 'Arduino Sketch\n/dev/ttyACM0'),
    (12.3, 4.35, 0.9, 'Pin 9\nSG90'),
]
for (x, y, w, txt) in labels:
    add_text(slide, txt, x, y, w, 0.7, size=11, color=LIGHT_GRAY, align=PP_ALIGN.CENTER)

# ─── SLIDE 3: How it works ─────────────────────────────────────────
slide = add_slide()
add_title_bar(slide, '⚙️ How It Works')

steps = [
    (BLUE,   '1', 'Webcam captures a frame',       'USB webcam plugged into KV260  →  /dev/video0  →  640×480 image'),
    (BLUE,   '2', 'Convert to grayscale',           'OpenCV converts color frame to grayscale for faster face detection'),
    (AMD_RED,'3', 'Run Haar Cascade detector',      'OpenCV built-in face model scans the frame — no GPU or DPU needed'),
    (GREEN,  '4', 'Face found?',                    'YES → send sweep command to Mega  |  NO → wait and scan again'),
    (ORANGE, '5', 'Send angles over USB serial',    'Python pyserial writes 0, 90, 180, 90, 0 to /dev/ttyACM0 at 9600 baud'),
    (YELLOW, '6', 'Servo sweeps',                   'Elegoo Mega reads angles → moves SG90 servo — repeats every 10 seconds'),
]

for i, (color, num, title, desc) in enumerate(steps):
    y = 1.4 + i * 0.85
    add_rect(slide, 0.4, y, 0.5, 0.6, color)
    add_text(slide, num,   0.4,  y+0.1,  0.5, 0.45, size=20, bold=True, color=DARK_BG, align=PP_ALIGN.CENTER)
    add_text(slide, title, 1.1,  y+0.05, 4.2, 0.4,  size=14, bold=True, color=color)
    add_text(slide, desc,  1.1,  y+0.42, 11.8,0.35, size=12, color=LIGHT_GRAY)

# ─── SLIDE 4: Hardware Wiring ──────────────────────────────────────
slide = add_slide()
add_title_bar(slide, '🔌 Hardware Wiring')

# Webcam box
add_rect(slide, 0.4, 1.5, 2.5, 1.0, BLUE)
add_text(slide, '📷 USB Webcam', 0.45, 1.6, 2.4, 0.7, size=14, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
add_text(slide, '↓ USB cable', 1.25, 2.55, 1.2, 0.4, size=13, color=LIGHT_GRAY, align=PP_ALIGN.CENTER)

# KV260 box
add_rect(slide, 0.4, 3.0, 2.5, 1.2, AMD_RED)
add_text(slide, '🧠 KV260\n192.168.68.60', 0.45, 3.1, 2.4, 1.0, size=14, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
add_text(slide, '↓ USB cable (/dev/ttyACM0)', 1.0, 4.25, 2.0, 0.4, size=11, color=LIGHT_GRAY, align=PP_ALIGN.CENTER)

# Mega box
add_rect(slide, 0.4, 4.7, 2.5, 1.0, ORANGE)
add_text(slide, '⚡ Elegoo Mega 2560', 0.45, 4.85, 2.4, 0.7, size=13, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

# Servo wiring
add_rect(slide, 4.0, 2.5, 5.5, 3.5, RGBColor(0x28,0x28,0x38))
add_text(slide, '🔌 Servo Wiring (SG90 → Mega)', 4.1, 2.6, 5.3, 0.5, size=15, bold=True, color=YELLOW)

wires = [
    (LIGHT_GRAY, 'Brown wire', '→', 'GND pin on Mega'),
    (AMD_RED,    'Red wire',   '→', '5V pin on Mega'),
    (ORANGE,     'Orange wire','→', 'Pin 9 on Mega'),
]
for i, (col, wire, arr, dest) in enumerate(wires):
    y = 3.3 + i * 0.8
    add_rect(slide, 4.2, y, 1.8, 0.55, col)
    add_text(slide, wire, 4.25, y+0.08, 1.7, 0.4, size=13, bold=True, color=DARK_BG, align=PP_ALIGN.CENTER)
    add_text(slide, arr,  6.1,  y+0.08, 0.4, 0.4, size=16, bold=True, color=WHITE)
    add_text(slide, dest, 6.55, y+0.08, 2.7, 0.4, size=13, color=WHITE)

# Servo box
add_rect(slide, 4.0, 5.6, 5.5, 1.0, GREEN)
add_text(slide, '🔄 SG90 Servo  (sweeps 0 → 90 → 180 → 90 → 0)', 4.1, 5.75, 5.3, 0.7, size=13, bold=True, color=DARK_BG, align=PP_ALIGN.CENTER)

# Key note
add_rect(slide, 9.8, 2.5, 3.2, 4.1, RGBColor(0x28,0x28,0x38))
add_text(slide, '💡 Key Notes', 9.9, 2.6, 3.0, 0.45, size=14, bold=True, color=YELLOW)
notes = [
    '• Serial baud: 9600',
    '• chmod 666 /dev/ttyACM0',
    '• vitis-env python3',
    '• Wait 3s after serial',
    '  open (Mega boot)',
    '• Stand 1-2m from',
    '  webcam for detection',
    '• Only 1 script can use',
    '  serial at a time',
]
for i, note in enumerate(notes):
    add_text(slide, note, 9.9, 3.1 + i*0.44, 3.0, 0.42, size=11, color=LIGHT_GRAY)

# ─── SLIDE 5: Demo Results ─────────────────────────────────────────
slide = add_slide()
add_title_bar(slide, '✅ Demo Results')

add_text(slide, 'What happens when you walk in front of the camera:', 0.5, 1.4, 12, 0.5, size=18, bold=True, color=WHITE)

results = [
    (GREEN,  '✅ Face detected',     '"Face detected! (1 face(s))" — printed every 0.5 seconds'),
    (ORANGE, '🔄 Servo sweeps',      '0° → 90° → 180° → 90° → 0°  —  smooth sweep in ~4 seconds'),
    (BLUE,   '⏱️ Every 10 seconds',  'While face is visible, servo re-sweeps every 10 seconds'),
    (YELLOW, '😶 No face → silent',  'When you walk away — servo stays still, camera keeps scanning'),
    (LIGHT_GRAY,'📊 Performance',    '~83ms per frame on ARM CPU  ·  ~12fps  ·  no GPU/DPU needed'),
]

for i, (color, title, desc) in enumerate(results):
    y = 2.1 + i * 0.95
    add_rect(slide, 0.4, y, 0.08, 0.65, color)
    add_text(slide, title, 0.65, y+0.1,  4.5, 0.5, size=15, bold=True, color=color)
    add_text(slide, desc,  5.3,  y+0.12, 7.7, 0.45, size=13, color=LIGHT_GRAY)

add_rect(slide, 0.4, 7.0, 12.4, 0.35, RGBColor(0x00,0x40,0x00))
add_text(slide, '🎉  Full pipeline working: Webcam → KV260 → Face Detection → Serial → Elegoo Mega → Servo Sweep',
         0.6, 7.03, 12.0, 0.3, size=12, bold=True, color=GREEN)

pptx_path = 'C:/Users/hemn/kv260/robotics/simple_project.pptx'
prs.save(pptx_path)
print('Saved:', pptx_path)
