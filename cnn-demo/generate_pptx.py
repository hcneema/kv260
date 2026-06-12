"""
Generate PowerPoint presentation for Kria KV260 bring-up session.
Run on Kria board: ~/vitis-env/bin/python3 generate_pptx.py
Then scp back to host.
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
import io, os

# AMD color palette
AMD_RED   = RGBColor(0xE8, 0x39, 0x1D)
AMD_DARK  = RGBColor(0x1A, 0x1A, 0x2E)
WHITE     = RGBColor(0xFF, 0xFF, 0xFF)
GRAY      = RGBColor(0xAA, 0xAA, 0xAA)
GREEN     = RGBColor(0x4C, 0xAF, 0x50)
YELLOW    = RGBColor(0xFF, 0xC1, 0x07)
BLUE      = RGBColor(0x21, 0x96, 0xF3)
DARK_BG   = RGBColor(0x0F, 0x0F, 0x1A)

prs = Presentation()
prs.slide_width  = Inches(13.33)
prs.slide_height = Inches(7.5)

BLANK = prs.slide_layouts[6]  # completely blank

def add_slide():
    return prs.slides.add_slide(BLANK)

def bg(slide, color=DARK_BG):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color

def box(slide, left, top, width, height, text="", fontsize=18,
        bold=False, color=WHITE, bgcolor=None, align=PP_ALIGN.LEFT,
        border_color=None):
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    if bgcolor:
        fill = txBox.fill
        fill.solid()
        fill.fore_color.rgb = bgcolor
    if border_color:
        txBox.line.color.rgb = border_color
        txBox.line.width = Pt(1)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(fontsize)
    run.font.bold = bold
    run.font.color.rgb = color
    return txBox

def accent_bar(slide, top):
    bar = slide.shapes.add_shape(1, Inches(0), Inches(top), Inches(0.07), Inches(0.6))
    bar.fill.solid()
    bar.fill.fore_color.rgb = AMD_RED
    bar.line.fill.background()

def rect(slide, left, top, width, height, color, alpha=None):
    shape = slide.shapes.add_shape(1, Inches(left), Inches(top), Inches(width), Inches(height))
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape

def heading(slide, text, top=0.4):
    accent_bar(slide, top)
    box(slide, 0.22, top, 12, 0.6, text, fontsize=28, bold=True, color=WHITE)

def bullet(slide, items, left, top, width=5.5, fontsize=15, color=WHITE):
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(5))
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        run = p.add_run()
        run.text = "  ▸  " + item
        run.font.size = Pt(fontsize)
        run.font.color.rgb = color
        p.space_after = Pt(6)

def stat_card(slide, left, top, number, label, color=AMD_RED):
    w, h = 2.8, 1.5
    r = rect(slide, left, top, w, h, RGBColor(0x16, 0x21, 0x3E))
    box(slide, left+0.1, top+0.1, w-0.2, 0.8, number, fontsize=32, bold=True, color=color, align=PP_ALIGN.CENTER)
    box(slide, left+0.1, top+0.9, w-0.2, 0.5, label.upper(), fontsize=10, color=GRAY, align=PP_ALIGN.CENTER)

def code_block(slide, left, top, width, height, code):
    r = rect(slide, left, top, width, height, RGBColor(0x1E, 0x1E, 0x2E))
    txBox = slide.shapes.add_textbox(Inches(left+0.15), Inches(top+0.1), Inches(width-0.3), Inches(height-0.2))
    tf = txBox.text_frame
    tf.word_wrap = False
    for i, line in enumerate(code.split('\n')):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        run = p.add_run()
        run.text = line
        run.font.size = Pt(11)
        run.font.name = "Cascadia Code"
        run.font.color.rgb = RGBColor(0x7E, 0xC8, 0xE3)

# ─────────────────────────────────────────
# SLIDE 1: Title
# ─────────────────────────────────────────
s = add_slide()
bg(s)
rect(s, 0, 0, 13.33, 7.5, RGBColor(0x1A, 0x0A, 0x0A))  # subtle tint
rect(s, 0, 0, 0.12, 7.5, AMD_RED)  # left accent stripe

box(s, 0.4, 0.8, 10, 0.5, "SESSION REPORT  ·  2026-06-07", fontsize=14, color=GRAY)
box(s, 0.4, 1.4, 11, 1.2, "Kria KV260 Bring-up", fontsize=46, bold=True, color=WHITE)
box(s, 0.4, 2.6, 11, 0.8, "& CNN Inference Pipeline", fontsize=36, bold=True, color=AMD_RED)
box(s, 0.4, 3.6, 9,  0.5,
    "End-to-end setup of the Xilinx AMD Kria KV260 Vision AI Starter Kit —\n"
    "from zero network access to running MobileNetV2 inference on-board.",
    fontsize=15, color=GRAY)

badges = [("SSH Connected","green"), ("CNN Running","green"), ("DPU — Next","yellow"), ("YOLO — Planned","blue")]
bx = 0.4
for label, col in badges:
    c = {"green": GREEN, "yellow": YELLOW, "blue": BLUE}[col]
    r = rect(s, bx, 4.5, 1.9, 0.4, RGBColor(0x1A,0x1A,0x2E))
    box(s, bx+0.05, 4.5, 1.8, 0.4, label, fontsize=12, color=c, align=PP_ALIGN.CENTER)
    bx += 2.05

box(s, 0.4, 6.8, 5, 0.4, "hemn  ·  AMD Engineering", fontsize=12, color=RGBColor(0x55,0x55,0x55))

# ─────────────────────────────────────────
# SLIDE 2: Board Overview
# ─────────────────────────────────────────
s = add_slide()
bg(s)
heading(s, "Board Overview")

# HW card
rect(s, 0.3, 1.3, 5.9, 5.6, RGBColor(0x1A,0x1A,0x2E))
box(s, 0.5, 1.4, 5, 0.4, "HARDWARE", fontsize=12, bold=True, color=AMD_RED)
bullet(s, [
    "Kria KV260 Vision AI Starter Kit",
    "SoC: Zynq UltraScale+ MPSoC",
    "CPU: Quad-core ARM Cortex-A53 @ 1.3GHz",
    "FPGA: PL fabric + AI Engine",
    "RAM: 4 GB LPDDR4",
    "Storage: 29 GB eMMC",
], 0.5, 1.9, 5.5, fontsize=14)

# SW card
rect(s, 6.8, 1.3, 5.9, 5.6, RGBColor(0x1A,0x1A,0x2E))
box(s, 7.0, 1.4, 5, 0.4, "SOFTWARE", fontsize=12, bold=True, color=AMD_RED)
bullet(s, [
    "Ubuntu 24.04.4 LTS (Noble Numbat)",
    "Kernel: 6.8.0-1024-xilinx (aarch64)",
    "Python 3.12.3",
    "xmutil + dfx-mgr (overlay mgmt)",
    "fpga-manager-xlnx",
    "libxaiengine 2024.2",
], 7.0, 1.9, 5.5, fontsize=14)

# warning bar
rect(s, 0.3, 6.5, 12.4, 0.6, RGBColor(0x33, 0x29, 0x00))
box(s, 0.5, 6.55, 12, 0.45,
    "⚠  Ubuntu 24.04 is cutting-edge for Kria — Vitis AI apt packages officially target Ubuntu 22.04 and below.",
    fontsize=12, color=YELLOW)

# ─────────────────────────────────────────
# SLIDE 3: Network Discovery
# ─────────────────────────────────────────
s = add_slide()
bg(s)
heading(s, "Challenge 1 — Network Discovery")

rect(s, 0.3, 1.3, 5.9, 5.6, RGBColor(0x1A,0x1A,0x2E))
box(s, 0.5, 1.4, 5, 0.3, "PROBLEMS", fontsize=12, bold=True, color=AMD_RED)
bullet(s, [
    "Host PC and Kria on different WiFi networks",
    "SSH (port 22) disabled by default",
    "mDNS kria.local didn't resolve",
    "Hostname xsjkria02x not in DNS",
    "Port scan: no open TCP ports found",
], 0.5, 1.8, 5.5, fontsize=13)

box(s, 0.5, 4.2, 5, 0.3, "RESOLUTION", fontsize=12, bold=True, color=GREEN)
bullet(s, [
    "Physical access: opened terminal on board",
    "Connected both devices to same WiFi",
    "Board IP confirmed: 192.168.68.60",
    "TTL=64 on ping → Linux confirmed",
], 0.5, 4.6, 5.5, fontsize=13, color=RGBColor(0xCC,0xFF,0xCC))

code_block(s, 6.8, 1.3, 5.9, 3.2,
"""# ARP scan on host
arp -a

# Wrong subnet initially:
192.168.1.116  64-ff-0a-85-f8-d2

# After switching WiFi:
ping 192.168.68.60
Reply: time=4ms TTL=64""")

code_block(s, 6.8, 4.7, 5.9, 2.2,
"""# Enable SSH on Kria
sudo systemctl enable --now ssh
ip addr show  # → 192.168.68.60""")

# ─────────────────────────────────────────
# SLIDE 4: SSH Setup
# ─────────────────────────────────────────
s = add_slide()
bg(s)
heading(s, "SSH Key Authentication Setup")

rect(s, 0.3, 1.3, 5.9, 3.5, RGBColor(0x1A,0x1A,0x2E))
box(s, 0.5, 1.4, 5, 0.3, "STEPS", fontsize=12, bold=True, color=AMD_RED)
bullet(s, [
    "Enabled SSH on board via physical terminal",
    "Generated ED25519 key pair on host PC",
    "Copied public key to board's authorized_keys",
    "Verified passwordless login",
], 0.5, 1.8, 5.5, fontsize=13)

for label, col, y in [("No password needed", GREEN, 5.1), ("Automated iteration enabled", GREEN, 5.6)]:
    r = rect(s, 0.5, y, 5.5, 0.4, RGBColor(0x1A,0x1A,0x2E))
    box(s, 0.6, y, 5.3, 0.4, "✓  " + label, fontsize=13, color=col)

code_block(s, 6.8, 1.3, 5.9, 5.8,
"""# Kria physical terminal
sudo systemctl enable --now ssh

# Host PC (PowerShell)
ssh-keygen -t ed25519 -f ~/.ssh/id_ed25519

type ~/.ssh/id_ed25519.pub | ssh \\
  ubuntu@192.168.68.60 \\
  "mkdir -p ~/.ssh && \\
   cat >> ~/.ssh/authorized_keys"

# Verify passwordless access
ssh ubuntu@192.168.68.60 "echo ok"
connected successfully ✓""")

# ─────────────────────────────────────────
# SLIDE 5: Software Stack
# ─────────────────────────────────────────
s = add_slide()
bg(s)
heading(s, "Software Stack Decision")

for i, (title, items, color, mark) in enumerate([
    ("Vitis AI (apt)", ["Not available for Ubuntu 24.04","Targets Ubuntu 20.04/22.04","No apt package found"], AMD_RED, "❌"),
    ("PYNQ",          ["pip install failed on Python 3.12","Build error on aarch64","Requires older Python"], AMD_RED, "❌"),
    ("ONNX Runtime",  ["Available for aarch64","Supports Python 3.12","Runs MobileNet/YOLO/ResNet","Extensible to DPU EP"], GREEN, "✅"),
]):
    left = 0.3 + i * 4.3
    border = GREEN if mark == "✅" else RGBColor(0x33,0x33,0x33)
    rect(s, left, 1.3, 4.0, 5.4, RGBColor(0x1A,0x1A,0x2E))
    box(s, left+0.15, 1.4, 3.7, 0.5, mark + "  " + title, fontsize=14, bold=True, color=color)
    bullet(s, items, left+0.15, 2.0, 3.6, fontsize=13,
           color=RGBColor(0xCC,0xFF,0xCC) if mark=="✅" else GRAY)

box(s, 0.3, 6.5, 12.5, 0.4, "Installation:  python3 -m venv ~/vitis-env  &&  pip install onnxruntime numpy pillow",
    fontsize=13, color=RGBColor(0x7E,0xC8,0xE3))

# ─────────────────────────────────────────
# SLIDE 6: CNN Pipeline
# ─────────────────────────────────────────
s = add_slide()
bg(s)
heading(s, "CNN Inference Pipeline")

rect(s, 0.3, 1.3, 5.9, 5.6, RGBColor(0x1A,0x1A,0x2E))
box(s, 0.5, 1.4, 5, 0.3, "MODEL", fontsize=12, bold=True, color=AMD_RED)
bullet(s, [
    "MobileNetV2 — ImageNet (1000 classes)",
    "Format: ONNX  |  Size: 14 MB",
    "Input: 224×224 RGB",
    "Source: ONNX Model Zoo (official)",
], 0.5, 1.8, 5.5, fontsize=13)

box(s, 0.5, 3.6, 5, 0.3, "PIPELINE", fontsize=12, bold=True, color=AMD_RED)
bullet(s, [
    "Load image → resize to 224×224",
    "Normalize (ImageNet mean/std)",
    "HWC → CHW → add batch dim",
    "ONNX Runtime inference session",
    "Softmax → Top-5 predictions",
], 0.5, 4.0, 5.5, fontsize=13)

code_block(s, 6.8, 1.3, 5.9, 5.8,
"""session = ort.InferenceSession(
    "mobilenet.onnx",
    providers=["CPUExecutionProvider"])

# Preprocess
img = Image.open("fish.jpg").resize((224,224))
x = np.array(img, dtype=np.float32)
x = (x/255.0 - mean) / std
x = x.transpose(2,0,1)[None]

# Infer
out = session.run(None, {"input": x})
top5 = np.argsort(out[0][0])[::-1][:5]""")

# ─────────────────────────────────────────
# SLIDE 7: Results (with fish image)
# ─────────────────────────────────────────
s = add_slide()
bg(s)
heading(s, "Results — MobileNetV2 on Kria KV260")

# Stats
stat_card(s, 0.3, 1.4, "87.5 ms", "Inference Time")
stat_card(s, 3.3, 1.4, "14 MB", "Model Size")
stat_card(s, 6.3, 1.4, "1000", "Classes")
stat_card(s, 9.3, 1.4, "✓", "Working", color=GREEN)

# Result table
rect(s, 0.3, 3.2, 6.0, 3.8, RGBColor(0x1A,0x1A,0x2E))
box(s, 0.5, 3.3, 5, 0.3, "RUN DETAILS", fontsize=12, bold=True, color=AMD_RED)
rows = [
    ("Model",     "MobileNetV2-12 (ONNX)"),
    ("Runtime",   "ONNX Runtime 1.26.0"),
    ("Device",    "ARM Cortex-A53 (CPU)"),
    ("Inference", "87.5 ms"),
    ("Top-1",     "tench  ✓  (score: 17.54)"),
    ("Status",    "End-to-end working ✓"),
]
for i, (k, v) in enumerate(rows):
    y = 3.75 + i*0.53
    box(s, 0.5, y, 2.5, 0.45, k, fontsize=13, color=GRAY)
    col = GREEN if "✓" in v else WHITE
    box(s, 3.0, y, 3.2, 0.45, v, fontsize=13, color=col)

# Fish image
if os.path.exists("fish.jpg"):
    s.shapes.add_picture("fish.jpg", Inches(6.6), Inches(3.0), Inches(3.2), Inches(2.4))
    box(s, 6.6, 5.5, 3.2, 0.4, "Input: tench fish (ImageNet sample)", fontsize=11, color=GRAY, align=PP_ALIGN.CENTER)

# Top predictions
rect(s, 10.1, 3.0, 2.9, 3.9, RGBColor(0x1A,0x1A,0x2E))
box(s, 10.2, 3.1, 2.5, 0.3, "TOP-5", fontsize=12, bold=True, color=AMD_RED)
preds = [("1. tench","17.54",GREEN),("2. snoek","11.21",WHITE),
         ("3. coho salmon","10.44",WHITE),("4. bolete","10.28",WHITE),("5. goldfish","9.78",WHITE)]
for i,(name,score,c) in enumerate(preds):
    y = 3.55 + i*0.65
    box(s, 10.2, y, 1.8, 0.5, name, fontsize=12, color=c)
    box(s, 11.8, y, 1.0, 0.5, score, fontsize=12, color=c, align=PP_ALIGN.RIGHT)

# ─────────────────────────────────────────
# SLIDE 8: Roadmap
# ─────────────────────────────────────────
s = add_slide()
bg(s)
heading(s, "Roadmap")

phases = [
    ("✅ Phase 1 — Board Bring-up & CPU Inference", "DONE",
     "SSH setup · ONNX Runtime · MobileNetV2 · 87.5ms on ARM Cortex-A53",
     GREEN, RGBColor(0x0A,0x2A,0x0A)),
    ("⚡ Phase 2 — DPU Acceleration", "NEXT",
     "Load DPU overlay via xmutil · ONNX DPU EP · Target: ~5–10ms (8–17× speedup)",
     YELLOW, RGBColor(0x2A,0x22,0x00)),
    ("🎯 Phase 3 — YOLO Object Detection", "PLANNED",
     "Quantize YOLOv8 INT8 with Vitis AI · Deploy on DPU · Live camera inference",
     BLUE, RGBColor(0x00,0x0A,0x2A)),
    ("🔬 Phase 4 — Custom HLS CNN", "FUTURE",
     "hls4ml: Keras → HLS C++ · Custom accelerator with Vitis HLS",
     RGBColor(0x9C,0x27,0xB0), RGBColor(0x1A,0x00,0x2A)),
]
for i, (title, tag, desc, color, bg_col) in enumerate(phases):
    y = 1.4 + i * 1.45
    rect(s, 0.3, y, 12.4, 1.2, bg_col)
    rect(s, 0.3, y, 0.08, 1.2, color)
    box(s, 0.55, y+0.1, 9.5, 0.45, title, fontsize=16, bold=True, color=color)
    box(s, 10.5, y+0.1, 2.0, 0.4, tag, fontsize=13, bold=True, color=color, align=PP_ALIGN.RIGHT)
    box(s, 0.55, y+0.6, 11.5, 0.45, desc, fontsize=13, color=GRAY)

# ─────────────────────────────────────────
# SLIDE 9: Key Learnings
# ─────────────────────────────────────────
s = add_slide()
bg(s)
heading(s, "Key Learnings")

rect(s, 0.3, 1.3, 5.9, 4.8, RGBColor(0x1A,0x1A,0x2E))
box(s, 0.5, 1.4, 5, 0.3, "TECHNICAL", fontsize=12, bold=True, color=AMD_RED)
bullet(s, [
    "Ubuntu 24.04 on Kria is bleeding-edge — Vitis AI apt packages not yet available for Noble",
    "ONNX Runtime is a solid fallback for aarch64 inference without full Vitis AI stack",
    "xmutil + dfx-mgr is the correct path for DPU overlay management",
    "TTL=64 on ping quickly confirms Linux target",
], 0.5, 1.8, 5.5, fontsize=13)

rect(s, 6.8, 1.3, 5.9, 4.8, RGBColor(0x1A,0x1A,0x2E))
box(s, 7.0, 1.4, 5, 0.3, "PROCESS", fontsize=12, bold=True, color=AMD_RED)
bullet(s, [
    "Always verify host and target on same subnet before debugging SSH",
    "SSH key auth essential for automated remote iteration",
    "Physical board access saved significant time for initial setup",
    "Watch disk space — eMMC at 64% before Docker/model installs",
], 7.0, 1.8, 5.5, fontsize=13)

rect(s, 0.3, 6.2, 12.4, 0.9, RGBColor(0x1A,0x1A,0x2E))
box(s, 0.5, 6.3, 12, 0.7,
    "Bottom line: Full CNN inference pipeline working on Kria KV260 in a single session.  "
    "Next: DPU acceleration → 87ms → <10ms,  then YOLO.",
    fontsize=14, color=WHITE)

# Save
prs.save("kria-kv260-slides.pptx")
print("Saved: kria-kv260-slides.pptx")
